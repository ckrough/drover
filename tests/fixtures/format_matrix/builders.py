"""Synthetic fixture builders for the format-coverage matrix.

Each builder writes a privacy-safe fixture for a single extension to
the target directory. Content is the same minimal payload across
formats so the test can assert non-empty extraction without depending
on per-format text variation:

    SAMPLE_TEXT = "Acme Corp Invoice 2025-11-15 amount 100"

Some extensions (legacy `.doc`, `.xls`, `.ppt`, `.tif` alias) cannot
be created with pure-Python libraries without a Microsoft-Office
toolchain. For those, we emit a marker file the test treats as
"unsupported by fixture builder" and reports as `SKIP` rather than
PASS / FAIL / REGRESS.
"""

from __future__ import annotations

import zipfile
from email.message import EmailMessage
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from pathlib import Path

SAMPLE_TEXT = "Acme Corp Invoice 2025-11-15 amount 100"


def build_pdf(target: Path) -> None:
    """A born-digital one-page PDF with a real text layer."""
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(target))
    c.drawString(72, 720, SAMPLE_TEXT)
    c.save()


def build_text(target: Path) -> None:
    target.write_text(SAMPLE_TEXT + "\n")


def build_markdown(target: Path) -> None:
    target.write_text(f"# Invoice\n\n{SAMPLE_TEXT}\n")


def build_html(target: Path) -> None:
    target.write_text(f"<html><body><h1>Invoice</h1><p>{SAMPLE_TEXT}</p></body></html>")


def build_csv(target: Path) -> None:
    target.write_text("vendor,date,amount\nAcme Corp,2025-11-15,100\n")


def build_tsv(target: Path) -> None:
    target.write_text("vendor\tdate\tamount\nAcme Corp\t2025-11-15\t100\n")


def build_image(target: Path, fmt: str) -> None:
    """Render `SAMPLE_TEXT` to an image in the requested format."""
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (1200, 300), color="white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=36)
    except OSError:
        font = ImageFont.load_default()
    draw.text((40, 100), SAMPLE_TEXT, fill="black", font=font)
    save_kwargs: dict[str, object] = {"format": fmt}
    img.save(target, **save_kwargs)


def build_docx(target: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Invoice", level=1)
    doc.add_paragraph(SAMPLE_TEXT)
    doc.save(target)


def build_xlsx(target: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["vendor", "date", "amount"])
    ws.append(["Acme Corp", "2025-11-15", 100])
    wb.save(target)


def build_pptx(target: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    if title is not None:
        title.text = "Invoice"
    body.text = SAMPLE_TEXT
    prs.save(target)


def build_eml(target: Path) -> None:
    msg = EmailMessage()
    msg["Subject"] = "Invoice"
    msg["From"] = "billing@acme.example"
    msg["To"] = "you@example.com"
    msg.set_content(SAMPLE_TEXT)
    target.write_bytes(bytes(msg))


def build_rtf(target: Path) -> None:
    target.write_text(
        r"{\rtf1\ansi\ansicpg1252\cocoartf2580 \pard\sa200\sl276\slmult1 "
        + SAMPLE_TEXT.replace("-", r"舑 ")
        + r" \par}"
    )


def build_epub(target: Path) -> None:
    """Minimal but valid EPUB 3 (single chapter)."""
    container_xml = (
        '<?xml version="1.0"?>\n'
        '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">\n'
        "  <rootfiles>\n"
        '    <rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>\n'
        "  </rootfiles>\n"
        "</container>\n"
    )
    content_opf = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<package xmlns="http://www.idpf.org/2007/opf" version="3.0" '
        'unique-identifier="bookid">\n'
        '  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">\n'
        '    <dc:identifier id="bookid">id-acme-2025</dc:identifier>\n'
        "    <dc:title>Invoice</dc:title>\n"
        "    <dc:language>en</dc:language>\n"
        "  </metadata>\n"
        "  <manifest>\n"
        '    <item id="ch1" href="ch1.xhtml" media-type="application/xhtml+xml"/>\n'
        "  </manifest>\n"
        "  <spine>\n"
        '    <itemref idref="ch1"/>\n'
        "  </spine>\n"
        "</package>\n"
    )
    chapter = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        "<!DOCTYPE html>\n"
        '<html xmlns="http://www.w3.org/1999/xhtml">\n'
        f"<head><title>Invoice</title></head>\n"
        f"<body><h1>Invoice</h1><p>{SAMPLE_TEXT}</p></body>\n"
        "</html>\n"
    )
    with zipfile.ZipFile(target, "w") as zf:
        # mimetype must be the first entry, stored uncompressed
        zf.writestr(
            zipfile.ZipInfo("mimetype"),
            "application/epub+zip",
            compress_type=zipfile.ZIP_STORED,
        )
        zf.writestr("META-INF/container.xml", container_xml)
        zf.writestr("OEBPS/content.opf", content_opf)
        zf.writestr("OEBPS/ch1.xhtml", chapter)


def build_odt(target: Path) -> None:
    """Minimal valid ODT (text/styles/manifest only)."""
    content_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        "<office:document-content "
        'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
        'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" '
        'office:version="1.2">\n'
        "  <office:body><office:text>\n"
        f"    <text:p>{SAMPLE_TEXT}</text:p>\n"
        "  </office:text></office:body>\n"
        "</office:document-content>\n"
    )
    manifest_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        "<manifest:manifest "
        'xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0" '
        'manifest:version="1.2">\n'
        '  <manifest:file-entry manifest:full-path="/" '
        'manifest:media-type="application/vnd.oasis.opendocument.text"/>\n'
        '  <manifest:file-entry manifest:full-path="content.xml" '
        'manifest:media-type="text/xml"/>\n'
        "</manifest:manifest>\n"
    )
    with zipfile.ZipFile(target, "w") as zf:
        zf.writestr(
            zipfile.ZipInfo("mimetype"),
            "application/vnd.oasis.opendocument.text",
            compress_type=zipfile.ZIP_STORED,
        )
        zf.writestr("content.xml", content_xml)
        zf.writestr("META-INF/manifest.xml", manifest_xml)


# Map extension -> builder. Legacy formats (.doc/.xls/.ppt) intentionally
# left unbuilt: neither pure-Python loader supports them without a
# Microsoft-Office toolchain (libreoffice / antiword), and they exist on
# the supported list because unstructured proxies through external tools
# in environments that have them.
EXT_BUILDERS: dict[str, object] = {
    ".pdf": build_pdf,
    ".txt": build_text,
    ".md": build_markdown,
    ".html": build_html,
    ".htm": build_html,
    ".csv": build_csv,
    ".tsv": build_tsv,
    ".png": lambda p: build_image(p, "PNG"),
    ".jpg": lambda p: build_image(p, "JPEG"),
    ".jpeg": lambda p: build_image(p, "JPEG"),
    ".gif": lambda p: build_image(p, "GIF"),
    ".bmp": lambda p: build_image(p, "BMP"),
    ".tiff": lambda p: build_image(p, "TIFF"),
    ".tif": lambda p: build_image(p, "TIFF"),
    ".docx": build_docx,
    ".xlsx": build_xlsx,
    ".pptx": build_pptx,
    ".eml": build_eml,
    ".epub": build_epub,
    ".odt": build_odt,
    ".rtf": build_rtf,
}


# Legacy / unsupported by the fixture builder.
UNSUPPORTED_BY_BUILDER: frozenset[str] = frozenset({".doc", ".xls", ".ppt"})


def build_all(target_dir: Path) -> dict[str, Path]:
    """Build a fixture for every supported extension; return ext -> path map."""
    target_dir.mkdir(parents=True, exist_ok=True)
    paths: dict[str, Path] = {}
    for ext, builder in EXT_BUILDERS.items():
        out = target_dir / f"sample{ext}"
        builder(out)  # type: ignore[operator]
        paths[ext] = out
    return paths
